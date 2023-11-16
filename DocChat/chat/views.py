from django.contrib.auth.decorators import login_required
from django.http import JsonResponse
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import login, logout, authenticate

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores.faiss import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.callbacks import get_openai_callback
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory

import os
import tempfile
from dotenv import load_dotenv

from PyPDF2 import PdfReader
from docx2txt import docx2txt
from pptx import Presentation

from chat.forms import UserFileForm, SendMessageForm
from chat.models import Chat, Message, UserFile

load_dotenv()


# Create your views here.
def homepage(request):
    return render(request, 'chat/home.html')


def get_file_text(file):
    text = None
    if file:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            if file.name.endswith('.pdf'):
                pdf_reader = PdfReader(temp_file.name)
                text = ''.join(
                    page.extract_text() for page in pdf_reader.pages)
            elif file.name.endswith('.txt'):
                with open(temp_file.name, 'r') as f:
                    text = f.read()
            elif file.name.endswith('.docx'):
                text = docx2txt.process(temp_file.name)
            elif file.name.endswith('.pptx'):
                prs = Presentation(temp_file.name)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = '\n'.join(text_runs)
    return text


def get_text_fragments(text):
    text_separator = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    fragments = text_separator.split_text(text)
    return fragments


def get_vector(text_fragments):
    api_key = os.getenv("OPENAI_API_KEY")
    investments = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=api_key)
    vector = FAISS.from_texts(text_fragments, investments)
    return vector


def get_conversation_fragments(vector):
    bot_gpt = ChatOpenAI(model="gpt-3.5-turbo")
    memory = ConversationBufferMemory(memory_key='chat_history',
                                      return_messages=True)
    conversation_fragments = ConversationalRetrievalChain.from_llm(
        llm=bot_gpt,
        retriever=vector.as_retriever(),
        memory=memory
    )
    return conversation_fragments


@login_required
def upload_file(request):
    template_name = 'chat/upload_file.html'
    user = request.user

    if request.method == 'POST':
        form = UserFileForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_document = request.FILES['file']

            _, file_extension = os.path.splitext(pdf_document.name)
            if file_extension.lower() not in [".pdf", ".txt", ".docx",
                                              ".pptx"]:
                return JsonResponse(
                    {'error': 'Only PDF, TXT, DOCX, PPTX files'}, status=400)

            if pdf_document.size > 50 * 1024 * 1024:
                return JsonResponse({'error': "File size exceeds 50 MB."},
                                    status=400)

            if UserFile.objects.filter(user=user,
                                       title=pdf_document.name).exists():
                return JsonResponse(
                    {'error': 'A file with the same name already exists.'},
                    status=400)

            user_file = UserFile(user=user, title=pdf_document.name)
            user_file.content = get_file_text(pdf_document)
            user_file.save()
            return redirect(
                'homepage')
    else:
        print("GET")
        form = UserFileForm()
    return render(request, template_name, {'form': form})


def create_chat(request):
    template_name = 'chat/chat.html'
    return redirect(
        'homepage')


@login_required
def send_message(request, chat_id=None):
    template_name = 'chat/chat.html'
    if request.method == 'POST':
        form = SendMessageForm(request.POST)
        if form.is_valid():
            sender = request.user
            message_text = request.POST.get('message')
            documents_content = UserFile.objects.filter(user=request.user).values_list('content', flat=True)
            text_chunks = []
            for doc in documents_content:
                text_chunks.extend(get_text_fragments(doc))

            knowledge_base = get_vector(text_chunks)
            conversation_chain = get_conversation_fragments(knowledge_base)

            with get_openai_callback() as cb:
                response = conversation_chain({'question': message_text})

            chat_response = response["answer"]
            print(chat_response)

            if chat_id:
                print(chat_id)
                chat = get_object_or_404(Chat, id=chat_id, user=request.user)
            else:
                chat = Chat.objects.create(user=request.user,
                                           title=message_text)
                print('create chat', chat)
                chat_id = chat.id
            Message.objects.create(chat=chat, sender=sender,
                                             message=message_text)
            Message.objects.create(chat=chat, sender=sender,
                                   message=chat_response, answer=True)
            return redirect('send_message_id', chat_id=chat_id)

    else:
        chats = Chat.objects.filter(user=request.user)
        chat_list = [{'id': chat.id, 'title': chat.title} for chat in chats]
        if chat_id:
            messages = get_messages(chat_id)
        else:
            messages = []
        form = SendMessageForm()
        return render(request, template_name,
                      {'chat_id': chat_id, 'form': form,
                       'message': messages, 'chats': chat_list})


@login_required
def get_chats(request):
    template_name = 'chat/chat_list.html'
    chats = Chat.objects.filter(user=request.user)
    chat_list = [{'id': chat.id, 'title': chat.title} for chat in chats]
    docs = UserFile.objects.filter(user=request.user)
    docs_list = [{'id': doc.id, 'title': doc.title} for doc in docs]
    return render(request, template_name, {'chats': chat_list, 'docs': docs_list})


def get_messages(chat_id):
    chat = get_object_or_404(Chat, id=chat_id)
    messages = Message.objects.filter(chat=chat)
    message_list = [
        {'sender': message.sender.username, 'message': message.message,
         'timestamp': message.timestamp, 'answer': message.answer} for message in messages]
    return message_list


@login_required
def delete_document(request, doc_id):
    doc = get_object_or_404(UserFile, id=doc_id, user=request.user)
    doc.delete()
    return redirect('get_chats')


@login_required
def delete_chat(request, chat_id):
    chat = get_object_or_404(Chat, id=chat_id, user=request.user)
    Message.objects.filter(chat=chat).delete()
    chat.delete()
    return redirect('get_chats')

